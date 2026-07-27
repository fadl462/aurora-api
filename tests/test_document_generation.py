"""
Tests app/document_generation.py and the /v1/generated-documents
endpoints. Like test_files.py, this verifies real file bytes come out
the other end — parsing the generated .pptx/.docx/.xlsx back with the
same libraries used to build them, not just checking status codes.
"""

import io
from types import SimpleNamespace
from unittest.mock import MagicMock, patch

import anthropic
import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation

from app.document_generation import _build_docx, _build_pptx, _build_xlsx, generate_document


# --- Builder tests: pure, deterministic, no model involved ---


def test_build_pptx_produces_real_readable_deck():
    content = {
        "title": "Q3 Review",
        "slides": [
            {"heading": "Highlights", "bullets": ["Revenue up 12%", "Two new clients"]},
            {"heading": "Risks", "bullets": ["Vendor delay"]},
        ],
    }
    file_bytes = _build_pptx(content)
    prs = Presentation(io.BytesIO(file_bytes))
    assert len(prs.slides) == 3  # title slide + 2 content slides
    assert prs.slides[0].shapes.title.text == "Q3 Review"
    assert prs.slides[1].shapes.title.text == "Highlights"


def test_build_docx_produces_real_readable_document():
    content = {
        "title": "Project Brief",
        "sections": [{"heading": "Overview", "paragraphs": ["This project does X.", "It matters because Y."]}],
    }
    file_bytes = _build_docx(content)
    doc = DocxDocument(io.BytesIO(file_bytes))
    full_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Project Brief" in full_text
    assert "Overview" in full_text
    assert "This project does X." in full_text


def test_build_xlsx_produces_real_readable_workbook():
    content = {
        "title": "Budget",
        "sheets": [{"name": "2026", "headers": ["Item", "Cost"], "rows": [["Rent", 1000], ["Supplies", 250]]}],
    }
    file_bytes = _build_xlsx(content)
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    ws = wb["2026"]
    assert [c.value for c in ws[1]] == ["Item", "Cost"]
    assert [c.value for c in ws[2]] == ["Rent", 1000]


def test_build_xlsx_with_no_sheets_still_produces_valid_file():
    file_bytes = _build_xlsx({"title": "Empty", "sheets": []})
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    assert len(wb.sheetnames) == 1  # fallback default sheet, never a workbook with zero sheets


# --- Full pipeline: no API key -> real, labeled placeholder file ---


def test_no_api_key_returns_real_placeholder_pptx(monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    result = generate_document("A deck about our Q3 results", "pptx")
    assert result["is_placeholder"] is True
    prs = Presentation(io.BytesIO(result["file_bytes"]))
    assert len(prs.slides) >= 1


def test_no_api_key_returns_real_placeholder_docx(monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    result = generate_document("A one-pager on our new product", "docx")
    assert result["is_placeholder"] is True
    doc = DocxDocument(io.BytesIO(result["file_bytes"]))
    assert len(doc.paragraphs) >= 1


def test_no_api_key_returns_real_placeholder_xlsx(monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    result = generate_document("A budget tracker", "xlsx")
    assert result["is_placeholder"] is True
    wb = openpyxl.load_workbook(io.BytesIO(result["file_bytes"]))
    assert len(wb.sheetnames) >= 1


# --- Full pipeline: real (mocked) model call ---


def test_real_model_call_produces_real_content(monkeypatch):
    monkeypatch.setenv("ANTHROPIC_API_KEY", "fake-key-for-test")
    fake_json = (
        '{"title": "Hiring Plan", "sections": '
        '[{"heading": "Roles", "paragraphs": ["We need two engineers."]}]}'
    )
    fake_response = SimpleNamespace(
        content=[SimpleNamespace(type="text", text=fake_json)],
    )

    with patch("app.document_generation.anthropic.Anthropic") as MockClient:
        MockClient.return_value.messages.create.return_value = fake_response
        result = generate_document("Plan our hiring for next quarter", "docx")

    assert result["is_placeholder"] is False
    assert result["title"] == "Hiring Plan"
    doc = DocxDocument(io.BytesIO(result["file_bytes"]))
    full_text = "\n".join(p.text for p in doc.paragraphs)
    assert "We need two engineers." in full_text


def test_model_response_wrapped_in_code_fences_still_parses(monkeypatch):
    monkeypatch.setenv("ANTHROPIC_API_KEY", "fake-key-for-test")
    fenced = '```json\n{"title": "Fenced", "sections": [{"heading": "H", "paragraphs": ["P"]}]}\n```'
    fake_response = SimpleNamespace(content=[SimpleNamespace(type="text", text=fenced)])

    with patch("app.document_generation.anthropic.Anthropic") as MockClient:
        MockClient.return_value.messages.create.return_value = fake_response
        result = generate_document("something", "docx")

    assert result["is_placeholder"] is False
    assert result["title"] == "Fenced"


def test_invalid_json_from_model_falls_back_to_placeholder(monkeypatch):
    monkeypatch.setenv("ANTHROPIC_API_KEY", "fake-key-for-test")
    fake_response = SimpleNamespace(content=[SimpleNamespace(type="text", text="not json at all")])

    with patch("app.document_generation.anthropic.Anthropic") as MockClient:
        MockClient.return_value.messages.create.return_value = fake_response
        result = generate_document("something", "pptx")

    assert result["is_placeholder"] is True


def test_api_status_error_falls_back_to_placeholder(monkeypatch):
    monkeypatch.setenv("ANTHROPIC_API_KEY", "fake-key-for-test")
    with patch("app.document_generation.anthropic.Anthropic") as MockClient:
        MockClient.return_value.messages.create.side_effect = anthropic.APIStatusError(
            message="bad request",
            response=MagicMock(status_code=400),
            body={"error": {"message": "some upstream reason"}},
        )
        result = generate_document("something", "xlsx")

    assert result["is_placeholder"] is True


# --- Router-level tests: real HTTP calls through the FastAPI app ---


def test_generate_endpoint_requires_auth(client):
    response = client.post("/v1/generated-documents", json={"prompt": "test", "format": "docx"})
    assert response.status_code == 401


def test_generate_endpoint_creates_real_downloadable_file(client, auth_headers, monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    headers = auth_headers()

    create_response = client.post(
        "/v1/generated-documents",
        json={"prompt": "A short brief about our product", "format": "docx"},
        headers=headers,
    )
    assert create_response.status_code == 201
    body = create_response.json()
    assert body["format"] == "docx"
    assert body["is_placeholder"] is True
    assert body["size_bytes"] > 0

    download_response = client.get(f"/v1/generated-documents/{body['id']}/download", headers=headers)
    assert download_response.status_code == 200
    assert download_response.headers["content-type"].startswith("application/vnd.openxmlformats")
    doc = DocxDocument(io.BytesIO(download_response.content))
    assert len(doc.paragraphs) >= 1


def test_list_generated_documents_only_returns_own(client, auth_headers, monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    headers_a = auth_headers("owner_a@example.com")
    headers_b = auth_headers("owner_b@example.com")

    client.post("/v1/generated-documents", json={"prompt": "doc for A", "format": "pptx"}, headers=headers_a)

    list_a = client.get("/v1/generated-documents", headers=headers_a)
    list_b = client.get("/v1/generated-documents", headers=headers_b)

    assert len(list_a.json()) == 1
    assert len(list_b.json()) == 0


def test_cannot_download_another_users_generated_document(client, auth_headers, monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    headers_a = auth_headers("owner_c@example.com")
    headers_b = auth_headers("owner_d@example.com")

    create_response = client.post(
        "/v1/generated-documents", json={"prompt": "private", "format": "xlsx"}, headers=headers_a
    )
    doc_id = create_response.json()["id"]

    response = client.get(f"/v1/generated-documents/{doc_id}/download", headers=headers_b)
    assert response.status_code == 404
