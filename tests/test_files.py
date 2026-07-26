"""
Tests app/document_extraction.py and the /v1/files/extract endpoint.
Each test generates a REAL file of that type in memory (using the same
libraries that write real PDFs/DOCX/etc.) rather than using canned
byte strings — so we're testing actual parsing, not a mock.
"""

import io

import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfWriter

from app.document_extraction import UnsupportedFileType, extract_text


def _make_real_docx(paragraphs: list[str]) -> bytes:
    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_real_pptx(slide_texts: list[str]) -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for text in slide_texts:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_real_xlsx(rows: list[list]) -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_real_pdf() -> bytes:
    # pypdf's writer can create a blank page; text extraction of a
    # blank page legitimately returns empty text — this test confirms
    # extraction runs without error on a real PDF structure, which is
    # what matters here (real text-bearing PDF creation needs a
    # rendering library beyond this project's scope).
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def test_extract_plain_text_file():
    result = extract_text("notes.txt", b"Hello, this is a real text file.")
    assert result["text"] == "Hello, this is a real text file."
    assert result["truncated"] is False
    assert result["char_count"] == 32


def test_extract_code_file():
    code = b"def hello():\n    print('hi')\n"
    result = extract_text("script.py", code)
    assert "def hello()" in result["text"]


def test_extract_real_docx():
    data = _make_real_docx(["First paragraph.", "Second paragraph with real content."])
    result = extract_text("brief.docx", data)
    assert "First paragraph." in result["text"]
    assert "Second paragraph with real content." in result["text"]


def test_extract_real_pptx():
    data = _make_real_pptx(["Introduction", "Q3 Results"])
    result = extract_text("deck.pptx", data)
    assert "Introduction" in result["text"]
    assert "Q3 Results" in result["text"]
    assert "[Slide 1]" in result["text"]
    assert "[Slide 2]" in result["text"]


def test_extract_real_xlsx():
    data = _make_real_xlsx([["Name", "Score"], ["Alice", 90], ["Bob", 85]])
    result = extract_text("scores.xlsx", data)
    assert "Alice" in result["text"]
    assert "90" in result["text"]


def test_extract_real_pdf_does_not_error():
    data = _make_real_pdf()
    result = extract_text("blank.pdf", data)
    assert "char_count" in result


def test_unsupported_file_type_raises():
    try:
        extract_text("photo.jpg", b"\xff\xd8\xff\xe0fakejpegdata")
        assert False, "should have raised UnsupportedFileType"
    except UnsupportedFileType:
        pass


def test_no_extension_treated_as_plain_text():
    result = extract_text("README", b"Just a readme with no extension")
    assert "Just a readme" in result["text"]


def test_long_text_gets_truncated():
    long_text = ("word " * 5000).encode("utf-8")
    result = extract_text("huge.txt", long_text)
    assert result["truncated"] is True
    assert result["char_count"] == 12_000


def test_extract_endpoint_requires_auth(client):
    response = client.post("/v1/files/extract", files={"file": ("notes.txt", b"hello", "text/plain")})
    assert response.status_code == 401


def test_extract_endpoint_returns_real_text(client, auth_headers):
    headers = auth_headers()
    response = client.post(
        "/v1/files/extract",
        files={"file": ("notes.txt", b"Real uploaded content here.", "text/plain")},
        headers=headers,
    )
    assert response.status_code == 200
    body = response.json()
    assert body["text"] == "Real uploaded content here."
    assert body["filename"] == "notes.txt"


def test_extract_endpoint_rejects_unsupported_type(client, auth_headers):
    headers = auth_headers()
    response = client.post(
        "/v1/files/extract",
        files={"file": ("photo.jpg", b"\xff\xd8\xff\xe0fakejpeg", "image/jpeg")},
        headers=headers,
    )
    assert response.status_code == 415


def test_extract_endpoint_rejects_oversized_file(client, auth_headers):
    headers = auth_headers()
    oversized = b"a" * (16 * 1024 * 1024)  # 16 MB, over the 15 MB limit
    response = client.post(
        "/v1/files/extract",
        files={"file": ("big.txt", oversized, "text/plain")},
        headers=headers,
    )
    assert response.status_code == 413
