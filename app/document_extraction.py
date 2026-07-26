"""
Real text extraction from uploaded files, so the assistant can actually
read what someone attaches instead of just seeing a filename (which is
all the composer's attach button did before this).

Deliberately text-extraction only, not layout/formatting-preserving —
good enough for the model to reason about content, not a document
viewer. Each file type gets its own extractor so adding a new type
later is additive, not a rewrite.
"""

import io

import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

MAX_EXTRACTED_CHARS = 12_000

TEXT_LIKE_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".css",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".cpp", ".h", ".cs",
    ".go", ".rb", ".php", ".sql", ".sh", ".rs", ".swift", ".kt",
}


class UnsupportedFileType(Exception):
    pass


def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    doc = DocxDocument(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        shape_texts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text]
        if shape_texts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(shape_texts))
    return "\n\n".join(slides_text)


def _extract_xlsx(data: bytes) -> str:
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    sheets_text = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append(" | ".join("" if cell is None else str(cell) for cell in row))
            if len(rows) >= 200:  # cap per-sheet rows so one huge sheet can't blow the budget alone
                rows.append("… (truncated)")
                break
        if rows:
            sheets_text.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(sheets_text)


def _extract_plain_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def extract_text(filename: str, data: bytes) -> dict:
    """Returns {text, truncated, char_count} or raises UnsupportedFileType."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == ".pdf":
        text = _extract_pdf(data)
    elif ext == ".docx":
        text = _extract_docx(data)
    elif ext == ".pptx":
        text = _extract_pptx(data)
    elif ext == ".xlsx":
        text = _extract_xlsx(data)
    elif ext in TEXT_LIKE_EXTENSIONS or ext == "":
        text = _extract_plain_text(data)
    else:
        raise UnsupportedFileType(f"'{ext}' files aren't supported yet.")

    truncated = len(text) > MAX_EXTRACTED_CHARS
    if truncated:
        text = text[:MAX_EXTRACTED_CHARS]

    return {"text": text, "truncated": truncated, "char_count": len(text)}
