"""
Real document generation — the write-side counterpart to
document_extraction.py. Aurora could read PDFs/Word/PowerPoint/Excel
but never produce one; this closes that gap.

Same honesty constraints as app/orchestration.py, deliberately:

1. When no ANTHROPIC_API_KEY is configured, this returns a real file
   built from clearly-labeled placeholder content — not a fake success,
   not a 500. The file structure/builders get fully exercised either
   way, which is also what makes them testable without an API key.

2. The model is never asked to "design a PowerPoint" in prose — it's
   asked for one specific JSON shape per format, parsed strictly. A
   model reply that doesn't parse falls back to the same placeholder
   path as a missing API key, rather than shipping a broken or
   half-written file.

3. Building the actual .pptx/.docx/.xlsx bytes is pure, deterministic
   code (python-pptx/python-docx/openpyxl) — the model's job is only to
   decide *content*, never to touch file structure directly.
"""

import io
import json
import os
from typing import Literal

import anthropic
from docx import Document as DocxDocument
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Pt

DocFormat = Literal["pptx", "docx", "xlsx"]

GENERATION_MODEL = "claude-sonnet-5"
GENERATION_MAX_TOKENS = 2048

PROMPT_MAX_CHARS = 4_000  # generation prompts are short asks, not document uploads

_SCHEMA_BY_FORMAT: dict[DocFormat, str] = {
    "pptx": (
        '{"title": "<deck title>", "slides": ['
        '{"heading": "<slide heading>", "bullets": ["<point>", "..."]}, "..."]}'
    ),
    "docx": (
        '{"title": "<document title>", "sections": ['
        '{"heading": "<section heading>", "paragraphs": ["<paragraph>", "..."]}, "..."]}'
    ),
    "xlsx": (
        '{"title": "<workbook title>", "sheets": ['
        '{"name": "<sheet name>", "headers": ["<col>", "..."], '
        '"rows": [["<cell>", "..."], "..."]}, "..."]}'
    ),
}


class GenerationError(Exception):
    """Raised only for programmer errors (bad format) — real-world
    failures (no key, bad model output, upstream error) degrade to the
    placeholder path instead of raising, matching orchestration.py."""


def _system_prompt(doc_format: DocFormat) -> str:
    schema = _SCHEMA_BY_FORMAT[doc_format]
    return (
        "You produce structured content for a document generator. Respond with "
        "ONLY valid JSON matching this exact shape, no prose, no markdown code "
        f"fences, nothing else: {schema}\n"
        "Keep it substantive but concrete: 4-8 slides/sections/rows depending on "
        "what the request actually calls for, not padded filler."
    )


def _placeholder_content(doc_format: DocFormat, prompt: str, reason: str) -> dict:
    note = (
        f"Placeholder content ({reason}). No real model call was made for this "
        "document — set ANTHROPIC_API_KEY in the environment to generate real "
        "content from your prompt. Your request was: "
        f"\"{prompt[:200]}\""
    )
    if doc_format == "pptx":
        return {"title": "Untitled deck (placeholder)", "slides": [{"heading": "Placeholder", "bullets": [note]}]}
    if doc_format == "docx":
        return {"title": "Untitled document (placeholder)", "sections": [{"heading": "Placeholder", "paragraphs": [note]}]}
    return {"title": "Untitled workbook (placeholder)", "sheets": [{"name": "Sheet1", "headers": ["Note"], "rows": [[note]]}]}


def _parse_model_json(raw_text: str) -> dict | None:
    text = raw_text.strip()
    # Models occasionally wrap JSON in fences despite instructions not to —
    # strip them rather than failing on an otherwise-valid response.
    if text.startswith("```"):
        text = text.strip("`")
        if text.startswith("json"):
            text = text[4:]
    try:
        parsed = json.loads(text.strip())
    except (json.JSONDecodeError, ValueError):
        return None
    return parsed if isinstance(parsed, dict) else None


def _generate_content(prompt: str, doc_format: DocFormat) -> tuple[dict, bool]:
    """Returns (content_dict, is_placeholder)."""
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        return _placeholder_content(doc_format, prompt, "no ANTHROPIC_API_KEY configured"), True

    try:
        client = anthropic.Anthropic(api_key=api_key)
        response = client.messages.create(
            model=os.environ.get("ANTHROPIC_MODEL", GENERATION_MODEL),
            max_tokens=GENERATION_MAX_TOKENS,
            system=_system_prompt(doc_format),
            messages=[{"role": "user", "content": prompt[:PROMPT_MAX_CHARS]}],
        )
        raw_text = "".join(block.text for block in response.content if block.type == "text")
        parsed = _parse_model_json(raw_text)
        if parsed is None:
            return _placeholder_content(doc_format, prompt, "model response wasn't valid JSON"), True
        return parsed, False
    except anthropic.APIStatusError as e:
        detail = "unknown reason"
        try:
            detail = e.response.json().get("error", {}).get("message", detail)
        except Exception:  # noqa: BLE001 — best-effort detail extraction only
            pass
        return _placeholder_content(doc_format, prompt, f"model call failed: HTTP {e.status_code} — {detail}"), True
    except anthropic.APIConnectionError:
        return _placeholder_content(doc_format, prompt, "could not reach the Anthropic API"), True
    except Exception as e:  # noqa: BLE001 — generation must never 500 on an upstream failure
        return _placeholder_content(doc_format, prompt, f"unexpected error calling the model: {type(e).__name__}"), True


def _build_pptx(content: dict) -> bytes:
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = content.get("title", "Untitled deck")

    for slide_data in content.get("slides", []):
        slide = prs.slides.add_slide(body_layout)
        slide.shapes.title.text = slide_data.get("heading", "")
        body = slide.placeholders[1].text_frame
        bullets = slide_data.get("bullets", [])
        if bullets:
            body.text = str(bullets[0])
            for bullet in bullets[1:]:
                p = body.add_paragraph()
                p.text = str(bullet)
                p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_docx(content: dict) -> bytes:
    doc = DocxDocument()
    doc.add_heading(content.get("title", "Untitled document"), level=0)

    for section in content.get("sections", []):
        heading = section.get("heading")
        if heading:
            doc.add_heading(str(heading), level=1)
        for paragraph in section.get("paragraphs", []):
            doc.add_paragraph(str(paragraph))

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_xlsx(content: dict) -> bytes:
    wb = Workbook()
    wb.remove(wb.active)  # default blank sheet — every sheet below is a real one

    sheets = content.get("sheets", [])
    if not sheets:
        sheets = [{"name": "Sheet1", "headers": [], "rows": []}]

    for sheet_data in sheets:
        name = str(sheet_data.get("name", "Sheet"))[:31]  # Excel's own sheet-name limit
        ws = wb.create_sheet(title=name)
        headers = sheet_data.get("headers", [])
        if headers:
            ws.append(headers)
            for cell in ws[1]:
                cell.font = Font(bold=True)
        for row in sheet_data.get("rows", []):
            ws.append(row)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


_BUILDERS = {"pptx": _build_pptx, "docx": _build_docx, "xlsx": _build_xlsx}

MIME_TYPES = {
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


def generate_document(prompt: str, doc_format: DocFormat) -> dict:
    """Full pipeline: prompt -> structured content -> real file bytes.

    Returns {file_bytes, title, is_placeholder} — never raises for
    real-world failures (missing key, bad model output, upstream
    errors); those all degrade to a real, clearly-labeled placeholder
    file so the caller always gets something downloadable.
    """
    if doc_format not in _BUILDERS:
        raise GenerationError(f"Unsupported format: {doc_format}")

    content, is_placeholder = _generate_content(prompt, doc_format)
    file_bytes = _BUILDERS[doc_format](content)
    title = content.get("title") or "Untitled"
    return {"file_bytes": file_bytes, "title": title, "is_placeholder": is_placeholder}
